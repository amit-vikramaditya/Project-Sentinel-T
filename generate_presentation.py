"""
generate_presentation.py — Project Sentinel-T Slide Deck Generator
====================================================================
Generates a full PowerPoint presentation using python-pptx.
Charts are created with matplotlib and embedded as images.

Usage:
    python generate_presentation.py

Output:
    Sentinel_T_Presentation.pptx
"""

import io
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour Palette ────────────────────────────────────────────────────────────
C_DARK_BLUE   = RGBColor(0x0D, 0x1B, 0x2A)   # near-black blue (slide bg)
C_MID_BLUE    = RGBColor(0x1B, 0x3A, 0x6B)   # section headers
C_ACCENT      = RGBColor(0x00, 0xB4, 0xD8)   # cyan accent
C_GREEN       = RGBColor(0x2E, 0xCC, 0x71)   # PHYSICAL / safe
C_RED         = RGBColor(0xE7, 0x4C, 0x3C)   # ANOMALY / attack
C_YELLOW      = RGBColor(0xF3, 0x9C, 0x12)   # WARMUP
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GREY  = RGBColor(0xEC, 0xF0, 0xF1)
C_DARK_GREY   = RGBColor(0x2C, 0x3E, 0x50)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────

def _rgb(r, g, b):
    return RGBColor(r, g, b)


def _add_bg(slide, colour: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, colour=C_WHITE,
                 align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = font_name
    return txBox


def _add_line(slide, colour: RGBColor, left, top, width, height=Pt(2)):
    """Thin horizontal rule."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()


def _fig_to_stream(fig):
    """Return a BytesIO PNG stream from a matplotlib figure."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _add_chart_image(slide, stream, left, top, width, height):
    slide.shapes.add_picture(stream, left, top, width, height)


def _slide_header(slide, title: str, subtitle: str = ""):
    """Common dark-themed header for content slides."""
    _add_bg(slide, C_DARK_BLUE)
    _add_line(slide, C_ACCENT, Inches(0.3), Inches(0.85), Inches(12.73))
    _add_textbox(slide, title,
                 Inches(0.3), Inches(0.1), Inches(12.73), Inches(0.75),
                 font_size=28, bold=True, colour=C_ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        _add_textbox(slide, subtitle,
                     Inches(0.3), Inches(0.85), Inches(12.73), Inches(0.4),
                     font_size=14, colour=C_LIGHT_GREY, align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────────────
# Chart generators (matplotlib → PNG stream)
# ─────────────────────────────────────────────────────────────────────────────

def _chart_clock_skew():
    """Cumulative clock skew: real ECU vs smart attacker."""
    np.random.seed(42)
    n = 1000
    base = 0.010

    # Real ECU
    thermal = np.sin(np.linspace(0, 4, n)) * 0.00002
    jitter, intervals_ecu = 0.0, []
    for i in range(n):
        jitter += -0.15 * jitter + 0.00001 * np.random.normal()
        intervals_ecu.append(base + jitter + thermal[i])
    skew_ecu = np.cumsum(np.array(intervals_ecu) - base)

    # Smart attacker
    skew_smart = np.cumsum(np.random.normal(0, 0.00005, n))

    fig, ax = plt.subplots(figsize=(8, 3.5), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    ax.plot(skew_ecu * 1e6,  color="#2ECC71", linewidth=2,
            label="Real ECU (thermal drift)")
    ax.plot(skew_smart * 1e6, color="#E74C3C", linewidth=1.5,
            linestyle="--", alpha=0.8, label="Smart Attacker (Gaussian)")
    ax.set_xlabel("Message count", color="#ECF0F1", fontsize=11)
    ax.set_ylabel("Cumulative skew (µs)", color="#ECF0F1", fontsize=11)
    ax.set_title("Cumulative Clock Skew: Physical Drift vs Statistical Noise",
                 color="#00B4D8", fontsize=13, fontweight="bold")
    ax.tick_params(colors="#ECF0F1")
    ax.spines[:].set_color("#1B3A6B")
    ax.legend(facecolor="#1B3A6B", labelcolor="#ECF0F1", fontsize=10)
    ax.grid(True, alpha=0.2, color="#ECF0F1")
    return _fig_to_stream(fig)


def _chart_residual_comparison():
    """Bar chart: mean residual for ECU vs attackers."""
    categories = ["Real ECU\n(cangen)", "Smart Attacker\n(Bash sleep)", "Threshold\n(200 µs)"]
    values     = [0.53, 575, 200]
    colours    = ["#2ECC71", "#E74C3C", "#F39C12"]

    fig, ax = plt.subplots(figsize=(7, 3.5), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    bars = ax.bar(categories, values, color=colours, edgecolor="#0D1B2A", linewidth=0)
    ax.set_ylabel("Mean Residual Error (µs)", color="#ECF0F1", fontsize=11)
    ax.set_title("Residual Error: ECU vs Attacker vs Threshold",
                 color="#00B4D8", fontsize=13, fontweight="bold")
    ax.tick_params(colors="#ECF0F1")
    ax.spines[:].set_color("#1B3A6B")
    ax.set_yscale("log")
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, val * 1.5,
                f"{val:.2f} µs", ha="center", color="#ECF0F1",
                fontsize=10, fontweight="bold")
    ax.grid(True, alpha=0.2, axis="y", color="#ECF0F1")
    return _fig_to_stream(fig)


def _chart_kalman_drift():
    """Kalman estimated drift over time for real ECU vs smart attacker."""
    np.random.seed(0)
    from drift_tracker import DriftTracker
    from sentinel_generator import SentinelGenerator

    gen = SentinelGenerator(num_samples=600)
    ecu   = gen.generate_real_ecu(receiver_jitter=5e-5)
    smart = gen.generate_smart_attacker(receiver_jitter=5e-5)

    t_ecu   = DriftTracker()
    t_smart = DriftTracker()
    _, d_ecu   = t_ecu.process_stream(ecu)
    _, d_smart = t_smart.process_stream(smart)

    fig, ax = plt.subplots(figsize=(8, 3.2), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    ax.plot(d_ecu   * 1e6, color="#2ECC71", linewidth=2, label="Real ECU drift")
    ax.plot(d_smart * 1e6, color="#E74C3C", linewidth=1.5,
            linestyle="--", alpha=0.8, label="Smart Attacker drift")
    ax.set_xlabel("Message count", color="#ECF0F1", fontsize=11)
    ax.set_ylabel("Estimated drift (µs/interval)", color="#ECF0F1", fontsize=11)
    ax.set_title("Kalman Filter: Estimated Clock Drift Rate",
                 color="#00B4D8", fontsize=13, fontweight="bold")
    ax.tick_params(colors="#ECF0F1")
    ax.spines[:].set_color("#1B3A6B")
    ax.legend(facecolor="#1B3A6B", labelcolor="#ECF0F1", fontsize=10)
    ax.grid(True, alpha=0.2, color="#ECF0F1")
    return _fig_to_stream(fig)


def _chart_detection_rate():
    """Horizontal bar chart: detection rate by attack type."""
    attacks = ["Injection\nAttack", "Fuzzing\nAttack", "Smart\nInjection", "Replay\nAttack"]
    rates   = [99.1, 98.3, 85.2, 72.4]
    colours = ["#E74C3C", "#E74C3C", "#E67E22", "#9B59B6"]

    fig, ax = plt.subplots(figsize=(7, 3.5), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    bars = ax.barh(attacks, rates, color=colours, edgecolor="#0D1B2A")
    ax.axvline(95, color="#F39C12", linestyle="--", linewidth=2, label="Target: 95%")
    ax.set_xlabel("Detection Rate (%)", color="#ECF0F1", fontsize=11)
    ax.set_title("Detection Rate by Attack Type", color="#00B4D8",
                 fontsize=13, fontweight="bold")
    ax.tick_params(colors="#ECF0F1")
    ax.spines[:].set_color("#1B3A6B")
    ax.set_xlim(60, 105)
    ax.legend(facecolor="#1B3A6B", labelcolor="#ECF0F1", fontsize=10)
    for bar, rate in zip(bars, rates):
        ax.text(rate + 0.8, bar.get_y() + bar.get_height()/2,
                f"{rate:.1f}%", va="center", color="#ECF0F1",
                fontsize=10, fontweight="bold")
    ax.grid(True, alpha=0.2, axis="x", color="#ECF0F1")
    return _fig_to_stream(fig)


def _chart_residual_distribution():
    """Histogram: residual distribution for normal vs attack."""
    np.random.seed(7)
    from drift_tracker import DriftTracker
    from sentinel_generator import SentinelGenerator

    gen = SentinelGenerator(num_samples=2000)
    ecu   = gen.generate_real_ecu(receiver_jitter=5e-5)
    smart = gen.generate_smart_attacker(receiver_jitter=5e-5)

    t_ecu   = DriftTracker()
    t_smart = DriftTracker()
    r_ecu,   _ = t_ecu.process_stream(ecu)
    r_smart, _ = t_smart.process_stream(smart)

    r_ecu_us   = np.abs(r_ecu[10:])   * 1e6
    r_smart_us = np.abs(r_smart[10:]) * 1e6

    fig, ax = plt.subplots(figsize=(8, 3.5), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    ax.hist(r_ecu_us,   bins=60, alpha=0.7, color="#2ECC71", density=True,
            label=f"Real ECU  (µ={r_ecu_us.mean():.1f} µs)")
    ax.hist(r_smart_us, bins=60, alpha=0.7, color="#E74C3C", density=True,
            label=f"Smart Attacker  (µ={r_smart_us.mean():.1f} µs)")
    ax.axvline(200, color="#F39C12", linewidth=2, linestyle="--",
               label="Threshold: 200 µs")
    ax.set_xlabel("Residual Error (µs)", color="#ECF0F1", fontsize=11)
    ax.set_ylabel("Density", color="#ECF0F1", fontsize=11)
    ax.set_title("Residual Distribution: ECU vs Smart Attacker",
                 color="#00B4D8", fontsize=13, fontweight="bold")
    ax.tick_params(colors="#ECF0F1")
    ax.spines[:].set_color("#1B3A6B")
    ax.legend(facecolor="#1B3A6B", labelcolor="#ECF0F1", fontsize=10)
    ax.set_xlim(0, 600)
    ax.grid(True, alpha=0.2, color="#ECF0F1")
    return _fig_to_stream(fig)


def _chart_architecture_diagram():
    """Simple architecture block diagram."""
    fig, ax = plt.subplots(figsize=(10, 3.8), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    ax.axis("off")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 4)

    boxes = [
        (0.3,  1.0, 1.8, 2.0, "#1B3A6B", "#00B4D8", "CAN Bus\n(vcan0 / real)"),
        (2.8,  1.0, 1.8, 2.0, "#1B3A6B", "#00B4D8", "Linux Kernel\nSO_TIMESTAMP"),
        (5.3,  1.0, 1.8, 2.0, "#1B3A6B", "#2ECC71", "Kalman Filter\n(DriftTracker)"),
        (7.8,  0.4, 1.8, 1.2, "#1B3A4B", "#2ECC71", "PHYSICAL\n[SAFE]"),
        (7.8,  2.4, 1.8, 1.2, "#4B1B1B", "#E74C3C", "ANOMALY\n[ALERT]"),
    ]
    for x, y, w, h, bg, border, label in boxes:
        rect = mpatches.FancyBboxPatch(
            (x, y), w, h,
            boxstyle="round,pad=0.1",
            facecolor=bg, edgecolor=border, linewidth=2
        )
        ax.add_patch(rect)
        ax.text(x + w/2, y + h/2, label, ha="center", va="center",
                color="#ECF0F1", fontsize=10, fontweight="bold",
                multialignment="center")

    # Arrows
    arrow_props = dict(arrowstyle="->", color="#00B4D8", lw=2)
    ax.annotate("", xy=(2.8, 2.0), xytext=(2.1, 2.0), arrowprops=arrow_props)
    ax.annotate("", xy=(5.3, 2.0), xytext=(4.6, 2.0), arrowprops=arrow_props)
    ax.annotate("", xy=(7.8, 1.0), xytext=(7.1, 1.7), arrowprops=arrow_props)
    ax.annotate("", xy=(7.8, 3.0), xytext=(7.1, 2.3), arrowprops=arrow_props)

    # Labels on arrows
    ax.text(2.45, 2.15, "frame +\ntimestamp", ha="center", va="bottom",
            color="#ECF0F1", fontsize=8)
    ax.text(4.95, 2.15, "us-precision\ninterval", ha="center", va="bottom",
            color="#ECF0F1", fontsize=8)

    ax.set_title("Sentinel-T Three-Layer Architecture",
                 color="#00B4D8", fontsize=14, fontweight="bold", pad=10)
    return _fig_to_stream(fig)


def _chart_snr_vs_jitter():
    """SNR vs OS jitter level."""
    jitter_us = [0, 50, 100, 200, 300, 500]
    snr_values = [8400, 11.4, 6.2, 3.8, 2.6, 1.9]

    fig, ax = plt.subplots(figsize=(7, 3.5), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    ax.plot(jitter_us, snr_values, color="#00B4D8", linewidth=3,
            marker="o", markersize=8, markerfacecolor="#2ECC71")
    ax.axhline(1.0, color="#E74C3C", linewidth=2, linestyle="--",
               label="Minimum viable SNR = 1.0")
    ax.set_xlabel("OS Scheduling Jitter (µs std)", color="#ECF0F1", fontsize=11)
    ax.set_ylabel("Signal-to-Noise Ratio (×)", color="#ECF0F1", fontsize=11)
    ax.set_title("Detection SNR vs Receiver OS Jitter (log scale)",
                 color="#00B4D8", fontsize=13, fontweight="bold")
    ax.set_yscale("log")
    ax.tick_params(colors="#ECF0F1")
    ax.spines[:].set_color("#1B3A6B")
    ax.legend(facecolor="#1B3A6B", labelcolor="#ECF0F1", fontsize=10)
    ax.grid(True, alpha=0.2, color="#ECF0F1", which="both")
    ax.fill_between(jitter_us, snr_values, 1.0,
                    where=[s > 1.0 for s in snr_values],
                    alpha=0.12, color="#2ECC71", label="Detectable region")
    return _fig_to_stream(fig)


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def slide_title(prs: Presentation):
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)
    _add_bg(slide, C_DARK_BLUE)

    # Decorative top bar
    _add_line(slide, C_ACCENT, 0, Inches(0), SLIDE_W, Inches(0.08))

    # Accent left bar
    bar = slide.shapes.add_shape(1, 0, Inches(0.08), Inches(0.35), Inches(7.42))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_MID_BLUE
    bar.line.fill.background()

    _add_textbox(slide, "PROJECT SENTINEL-T",
                 Inches(0.5), Inches(1.2), Inches(12.33), Inches(1.3),
                 font_size=52, bold=True, colour=C_ACCENT, align=PP_ALIGN.CENTER)

    _add_textbox(slide,
                 "A Physical-Layer Intrusion Detection System\nfor Automotive CAN Networks",
                 Inches(0.5), Inches(2.7), Inches(12.33), Inches(1.0),
                 font_size=22, colour=C_WHITE, align=PP_ALIGN.CENTER)

    _add_line(slide, C_ACCENT, Inches(2.5), Inches(3.9), Inches(8.33))

    _add_textbox(slide, "Amit Vikramaditya  |  Automotive Software Engineering",
                 Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.5),
                 font_size=16, colour=C_LIGHT_GREY, align=PP_ALIGN.CENTER)

    _add_textbox(slide, "Detecting CAN Bus attacks by fingerprinting\nphysical ECU clock drift using Kernel-Level Timestamping",
                 Inches(0.5), Inches(4.7), Inches(12.33), Inches(0.9),
                 font_size=14, colour=_rgb(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

    # Bottom decorative bar
    _add_line(slide, C_ACCENT, 0, Inches(7.42), SLIDE_W)


def slide_agenda(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _slide_header(slide, "Agenda", "What we will cover today")

    items = [
        ("01", "Introduction", "CAN bus, attack surface, why security is hard"),
        ("02", "Existing Systems", "Survey of current CAN IDS approaches & limitations"),
        ("03", "Proposed Solution", "Physical-layer clock fingerprinting concept"),
        ("04", "System Architecture", "Three-layer pipeline + Kalman Filter"),
        ("05", "Results", "Live validation, datasets, stress tests"),
        ("06", "Conclusion", "Key findings, limitations, future work"),
    ]

    for i, (num, title, desc) in enumerate(items):
        row = i % 3
        col = i // 3
        x = Inches(0.4 + col * 6.5)
        y = Inches(1.2 + row * 1.85)

        box = slide.shapes.add_shape(1, x, y, Inches(6.0), Inches(1.6))
        box.fill.solid(); box.fill.fore_color.rgb = C_MID_BLUE
        box.line.color.rgb = C_ACCENT

        _add_textbox(slide, num,  x + Inches(0.1), y + Inches(0.05),
                     Inches(0.7), Inches(0.5),
                     font_size=22, bold=True, colour=C_ACCENT)
        _add_textbox(slide, title, x + Inches(0.8), y + Inches(0.05),
                     Inches(5.0), Inches(0.55),
                     font_size=17, bold=True, colour=C_WHITE)
        _add_textbox(slide, desc,  x + Inches(0.8), y + Inches(0.65),
                     Inches(5.0), Inches(0.8),
                     font_size=12, colour=C_LIGHT_GREY)


def slide_intro(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _slide_header(slide, "Introduction: The CAN Bus Security Problem",
                  "ISO 11898 | Designed in 1986 | No authentication | ~1.4 billion vehicles")

    _add_textbox(slide, "CAN Bus: The Nervous System of Every Car",
                 Inches(0.3), Inches(1.05), Inches(8), Inches(0.5),
                 font_size=16, bold=True, colour=C_ACCENT)

    facts = (
        "• 70–100 ECUs in a modern vehicle, all sharing one broadcast bus\n"
        "• No sender address · No encryption · No authentication\n"
        "• Every node receives every message — implicit trust\n"
        "• 2015: Jeep Cherokee remotely compromised via cellular → steering & brakes disabled"
    )
    _add_textbox(slide, facts,
                 Inches(0.3), Inches(1.6), Inches(7.8), Inches(2.2),
                 font_size=14, colour=C_WHITE)

    _add_textbox(slide, "Why Not Just Add Encryption?",
                 Inches(0.3), Inches(3.9), Inches(8), Inches(0.5),
                 font_size=16, bold=True, colour=C_ACCENT)

    blockers = (
        "① Only 8 bytes payload — HMAC alone is 16 bytes\n"
        "② 8/16-bit MCUs in legacy ECUs — no crypto hardware\n"
        "③ 1–10 ms timing budget — crypto overhead violates real-time constraints\n"
        "④ Key management across 100 ECUs from 50 suppliers — unsolved problem"
    )
    _add_textbox(slide, blockers,
                 Inches(0.3), Inches(4.45), Inches(7.8), Inches(2.0),
                 font_size=13, colour=C_WHITE)

    # Side highlight box
    box = slide.shapes.add_shape(1, Inches(8.6), Inches(1.05), Inches(4.4), Inches(5.4))
    box.fill.solid(); box.fill.fore_color.rgb = C_MID_BLUE
    box.line.color.rgb = C_ACCENT

    _add_textbox(slide, "Attack Surfaces",
                 Inches(8.7), Inches(1.15), Inches(4.2), Inches(0.5),
                 font_size=15, bold=True, colour=C_ACCENT)

    surfaces = (
        "🔌  OBD-II Port\n"
        "    Physical under dashboard\n\n"
        "📡  TCU (Telematics)\n"
        "    Cellular remote access\n\n"
        "🎵  Infotainment\n"
        "    Wi-Fi · Bluetooth · USB\n\n"
        "🚗  V2X Module\n"
        "    Over-the-air injection\n\n"
        "🔧  Aftermarket Dongles\n"
        "    Malicious OBD-II accessories"
    )
    _add_textbox(slide, surfaces,
                 Inches(8.7), Inches(1.7), Inches(4.1), Inches(4.5),
                 font_size=12, colour=C_WHITE)


def slide_existing_systems(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _slide_header(slide, "Existing Systems",
                  "Survey of CAN Intrusion Detection approaches and their limitations")

    cols = [
        ("Content / Payload\nIDS", "#1B3A6B",
         "Inspects data bytes for\nabnormal values",
         "❌ Defeated by data-\n    aware attackers\n❌ Cannot detect\n    masquerade attacks"),
        ("Frequency /\nTiming IDS", "#1B3A6B",
         "Monitors message rate\nfor injection patterns",
         "❌ Time-aligned attackers\n    bypass detection\n❌ High false-positive\n    rate on load spikes"),
        ("AUTOSAR SecOC\n(Cryptographic)", "#1B3A4B",
         "Adds HMAC to each frame;\nrequires key management",
         "❌ Requires new MCUs\n❌ Not backwards compat.\n❌ Key mgmt complexity"),
        ("Cho & Shin 2016\n(Clock Fingerprint)", "#1B3A4B",
         "Linear regression on\nclock skew, USENIX Sec.",
         "❌ No jitter separation\n❌ No state model\n❌ Clean environment only"),
        ("Sentinel-T\n(This Project) ✅", "#0D2B1B",
         "Kalman Filter separates\nphysical drift from OS jitter",
         "✅ Works on legacy ECUs\n✅ Real-time < 0.04 ms\n✅ Separates jitter"),
    ]

    for i, (title, bg, how, limits) in enumerate(cols):
        x = Inches(0.25 + i * 2.55)
        # Header
        hbox = slide.shapes.add_shape(1, x, Inches(1.15), Inches(2.45), Inches(0.65))
        hbox.fill.solid(); hbox.fill.fore_color.rgb = RGBColor(0x00, 0xB4, 0xD8)
        hbox.line.fill.background()
        _add_textbox(slide, title, x + Inches(0.05), Inches(1.18),
                     Inches(2.35), Inches(0.6),
                     font_size=11, bold=True, colour=C_DARK_BLUE, align=PP_ALIGN.CENTER)

        # Body
        bdy = slide.shapes.add_shape(1, x, Inches(1.8), Inches(2.45), Inches(5.3))
        bdy.fill.solid(); bdy.fill.fore_color.rgb = RGBColor(*[int(bg[i:i+2], 16) for i in (1, 3, 5)])
        bdy.line.color.rgb = C_ACCENT

        _add_textbox(slide, "How it works:", x + Inches(0.1), Inches(1.9),
                     Inches(2.25), Inches(0.35),
                     font_size=11, bold=True, colour=C_ACCENT)
        _add_textbox(slide, how, x + Inches(0.1), Inches(2.25),
                     Inches(2.25), Inches(1.2),
                     font_size=11, colour=C_WHITE)

        _add_textbox(slide, "Limitations:", x + Inches(0.1), Inches(3.55),
                     Inches(2.25), Inches(0.35),
                     font_size=11, bold=True, colour=C_YELLOW)
        _add_textbox(slide, limits, x + Inches(0.1), Inches(3.9),
                     Inches(2.25), Inches(2.0),
                     font_size=10, colour=C_WHITE)


def slide_proposed_solution(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _slide_header(slide, "Proposed Solution: Physical-Layer Clock Fingerprinting",
                  "Authenticate ECUs by what their silicon IS, not what they SAY")

    _add_textbox(slide, "The Core Insight",
                 Inches(0.3), Inches(1.1), Inches(12), Inches(0.45),
                 font_size=16, bold=True, colour=C_ACCENT)

    _add_textbox(slide,
                 "Every ECU contains a quartz crystal oscillator. "
                 "Its resonant frequency drifts slowly with temperature — "
                 "this drift is physically unique per device. "
                 "No software running on any OS can replicate it without measuring the actual crystal.",
                 Inches(0.3), Inches(1.55), Inches(12), Inches(0.8),
                 font_size=13, colour=C_WHITE)

    # Two-column: Real ECU vs Attacker
    for i, (label, colour, detail) in enumerate([
        ("Real ECU (Physical Crystal)", C_GREEN,
         "• Thermal drift: slow sinusoidal, correlated\n"
         "• O-U jitter: mean-reverting, bounded\n"
         "• Consecutive intervals are NOT independent\n"
         "• Drift rate = ±20 ppm over temperature range\n"
         "• Detectable with Kalman Filter"),
        ("Software Attacker (OS Timer)", C_RED,
         "• sleep(0.01) scheduling: memoryless\n"
         "• Each interval is independent of the last\n"
         "• Gaussian noise from scheduler quantisation\n"
         "• No correlation structure\n"
         "• Filter residual stays persistently high"),
    ]):
        x = Inches(0.3 + i * 6.55)
        hbox = slide.shapes.add_shape(1, x, Inches(2.5), Inches(6.2), Inches(0.5))
        hbox.fill.solid(); hbox.fill.fore_color.rgb = colour
        hbox.line.fill.background()
        _add_textbox(slide, label, x + Inches(0.1), Inches(2.52),
                     Inches(6.0), Inches(0.46),
                     font_size=13, bold=True, colour=C_DARK_BLUE)

        bbox = slide.shapes.add_shape(1, x, Inches(3.0), Inches(6.2), Inches(2.8))
        bbox.fill.solid(); bbox.fill.fore_color.rgb = C_MID_BLUE
        bbox.line.color.rgb = colour
        _add_textbox(slide, detail, x + Inches(0.1), Inches(3.05),
                     Inches(6.0), Inches(2.7), font_size=12, colour=C_WHITE)

    # Chart
    stream = _chart_clock_skew()
    _add_chart_image(slide, stream, Inches(0.3), Inches(5.85), Inches(12.7), Inches(1.5))


def slide_architecture(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _slide_header(slide, "System Architecture",
                  "Three-layer pipeline: Physical Bus → Kernel Timestamp Tap → Chronomorphic Engine")

    stream = _chart_architecture_diagram()
    _add_chart_image(slide, stream, Inches(0.3), Inches(1.1), Inches(12.7), Inches(3.0))

    # Kalman Filter box
    kf_box = slide.shapes.add_shape(1, Inches(0.3), Inches(4.2), Inches(12.7), Inches(2.9))
    kf_box.fill.solid(); kf_box.fill.fore_color.rgb = C_MID_BLUE
    kf_box.line.color.rgb = C_ACCENT

    _add_textbox(slide, "Phase-Velocity Kalman Filter  |  State: x = [φ (offset),  φ̇ (drift rate)]",
                 Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.45),
                 font_size=14, bold=True, colour=C_ACCENT)

    kf_text = (
        "  Predict:     x_pred = F · x                  P_pred = F · P · Fᵀ + Q          Q = 1×10⁻¹²  (slow drift)\n"
        "  Measure:    z = Δt − T_nominal                residual r = z − H · x_pred\n"
        "  Update:      K = P_pred · Hᵀ · (H · P_pred · Hᵀ + R)⁻¹        x = x_pred + K · r       R = 1×10⁻¹⁰  (OS jitter)\n\n"
        "  |r| < 200 µs  →  PHYSICAL  ✅           |r| ≥ 200 µs  →  ANOMALY  🚨           count < 10  →  WARMUP  ⏳"
    )
    _add_textbox(slide, kf_text,
                 Inches(0.4), Inches(4.8), Inches(12.5), Inches(2.0),
                 font_size=12, colour=C_WHITE, font_name="Courier New")


def slide_results_live(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _slide_header(slide, "Results: Live Monitor (Azure VM)",
                  "Ubuntu 24.04 + linux-modules-extra + vcan0 virtual CAN interface")

    # KPIs
    kpis = [
        ("0.53 µs", "ECU\nResidual"),
        ("350–800 µs", "Attacker\nResidual"),
        ("1.88×", "SNR"),
        ("< 0.04 ms", "Latency"),
        ("0%", "False\nPositive"),
        ("0%", "False\nNegative"),
    ]
    kpi_colours = [C_GREEN, C_RED, C_ACCENT, C_ACCENT, C_GREEN, C_GREEN]

    for i, ((val, label), colour) in enumerate(zip(kpis, kpi_colours)):
        x = Inches(0.3 + i * 2.12)
        box = slide.shapes.add_shape(1, x, Inches(1.1), Inches(2.0), Inches(1.5))
        box.fill.solid(); box.fill.fore_color.rgb = C_MID_BLUE
        box.line.color.rgb = colour
        _add_textbox(slide, val, x, Inches(1.15), Inches(2.0), Inches(0.8),
                     font_size=20, bold=True, colour=colour, align=PP_ALIGN.CENTER)
        _add_textbox(slide, label, x, Inches(1.95), Inches(2.0), Inches(0.5),
                     font_size=11, colour=C_LIGHT_GREY, align=PP_ALIGN.CENTER)

    stream = _chart_residual_comparison()
    _add_chart_image(slide, stream, Inches(0.3), Inches(2.8), Inches(6.7), Inches(4.5))

    stream2 = _chart_kalman_drift()
    _add_chart_image(slide, stream2, Inches(7.1), Inches(2.8), Inches(6.0), Inches(4.5))


def slide_results_datasets(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _slide_header(slide, "Results: Dataset Benchmark Validation",
                  "5 benchmark datasets · 6 ECU types · 4 attack scenarios")

    stream = _chart_detection_rate()
    _add_chart_image(slide, stream, Inches(0.3), Inches(1.1), Inches(6.7), Inches(4.0))

    stream2 = _chart_residual_distribution()
    _add_chart_image(slide, stream2, Inches(7.1), Inches(1.1), Inches(6.0), Inches(4.0))

    # Table
    table_data = [
        ["Dataset",           "Messages", "Detection Rate", "FPR"],
        ["Normal (baseline)", "5,400",    "—",              "~0%"],
        ["Injection Attack",  "15,700",   "99.1%",          "0.5%"],
        ["Smart Injection",   "15,700",   "85.2%",          "1.0%"],
        ["Fuzzing Attack",    "67,000",   "98.3%",          "0.3%"],
        ["Replay Attack",     "15,700",   "72.4%",          "1.2%"],
    ]

    rows = len(table_data)
    cols = len(table_data[0])
    tbl = slide.shapes.add_table(rows, cols,
                                  Inches(0.3), Inches(5.15),
                                  Inches(12.7), Inches(2.2)).table
    tbl.columns[0].width = Inches(3.5)
    tbl.columns[1].width = Inches(2.0)
    tbl.columns[2].width = Inches(3.7)
    tbl.columns[3].width = Inches(3.5)

    for r, row_data in enumerate(table_data):
        for c, cell_text in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = cell_text
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.name = "Calibri"
            run.font.size = Pt(12)
            run.font.bold = (r == 0)
            run.font.color.rgb = C_DARK_BLUE if r == 0 else C_WHITE
            fill = cell.fill
            fill.solid()
            if r == 0:
                fill.fore_color.rgb = C_ACCENT
            elif r % 2 == 0:
                fill.fore_color.rgb = C_MID_BLUE
            else:
                fill.fore_color.rgb = C_DARK_GREY


def slide_results_stress(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _slide_header(slide, "Results: OS Jitter Robustness",
                  "Detection SNR remains > 1.0 even at 500 µs OS scheduling jitter")

    stream = _chart_snr_vs_jitter()
    _add_chart_image(slide, stream, Inches(0.3), Inches(1.1), Inches(7.5), Inches(5.0))

    _add_textbox(slide, "Key Finding",
                 Inches(8.1), Inches(1.1), Inches(4.9), Inches(0.5),
                 font_size=16, bold=True, colour=C_ACCENT)

    findings = (
        "Even at 500 µs of OS scheduling\n"
        "jitter — typical of a heavily loaded\n"
        "non-RT Linux kernel — the SNR\n"
        "remains above 1.0.\n\n"
        "The physical drift signal is\n"
        "recoverable across all realistic\n"
        "operating conditions.\n\n"
        "The Kalman Filter's high R/Q ratio\n"
        "(= 100) is the key: it deliberately\n"
        "ignores single-measurement noise and\n"
        "tracks only persistent, structured\n"
        "deviations — exactly the signature\n"
        "that distinguishes silicon from software."
    )
    _add_textbox(slide, findings,
                 Inches(8.1), Inches(1.7), Inches(4.9), Inches(4.5),
                 font_size=13, colour=C_WHITE)


def slide_conclusion(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _slide_header(slide, "Conclusion & Future Work")

    findings = [
        ("✅  Kernel timestamping works",
         "SO_TIMESTAMP reduces noise floor by 1–2 orders vs user-space time.time()"),
        ("✅  Kalman Filter is effective",
         "Simple 2D state model separates thermal drift from scheduling jitter (SNR 1.88×)"),
        ("✅  Clear detection boundary",
         "0.53 µs (ECU) vs 350–800 µs (attacker) → 200 µs threshold, 0% FP/FN observed"),
        ("✅  Real-time capable",
         "< 0.04 ms per packet — < 0.4% of the 10 ms CAN cycle budget"),
        ("⚠️  Replay is hardest",
         "72% detection — similar timing to real ECU; requires multi-feature fusion"),
    ]

    for i, (heading, detail) in enumerate(findings):
        y = Inches(1.15 + i * 1.2)
        box = slide.shapes.add_shape(1, Inches(0.3), y, Inches(12.7), Inches(1.05))
        box.fill.solid()
        box.fill.fore_color.rgb = C_MID_BLUE if i < 4 else RGBColor(0x4B, 0x2A, 0x00)
        box.line.color.rgb = C_GREEN if i < 4 else C_YELLOW
        _add_textbox(slide, heading, Inches(0.4), y + Inches(0.05),
                     Inches(12.5), Inches(0.45),
                     font_size=14, bold=True,
                     colour=C_GREEN if i < 4 else C_YELLOW)
        _add_textbox(slide, detail, Inches(0.4), y + Inches(0.5),
                     Inches(12.5), Inches(0.45),
                     font_size=12, colour=C_WHITE)

    # Future work bar
    fw = slide.shapes.add_shape(1, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3))
    fw.fill.solid(); fw.fill.fore_color.rgb = C_DARK_GREY
    fw.line.fill.background()
    _add_textbox(slide,
                 "Future:  Phase 2 → Physical hardware (RPi + MCP2515)  |  "
                 "Phase 3 → Mamba SSM  |  Phase 4 → Multi-feature (clock + voltage + sequence)",
                 Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.26),
                 font_size=11, colour=C_ACCENT)


def slide_thankyou(prs: Presentation):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    _add_bg(slide, C_DARK_BLUE)
    _add_line(slide, C_ACCENT, 0, 0, SLIDE_W)
    _add_line(slide, C_ACCENT, 0, Inches(7.42), SLIDE_W)

    _add_textbox(slide, "Thank You",
                 Inches(0), Inches(2.0), Inches(13.33), Inches(1.5),
                 font_size=64, bold=True, colour=C_ACCENT, align=PP_ALIGN.CENTER)

    _add_textbox(slide, "Project Sentinel-T  |  Automotive Software Engineering",
                 Inches(0), Inches(3.7), Inches(13.33), Inches(0.6),
                 font_size=20, colour=C_WHITE, align=PP_ALIGN.CENTER)

    _add_textbox(slide, "github.com/amit-vikramaditya/Project-Sentinel-T",
                 Inches(0), Inches(4.4), Inches(13.33), Inches(0.5),
                 font_size=16, colour=_rgb(0x00, 0xB4, 0xD8), align=PP_ALIGN.CENTER)

    _add_textbox(slide,
                 "References: Miller & Valasek (2015) · Cho & Shin, USENIX Security (2016)\n"
                 "Müter & Asaj, IEEE IV (2011) · ISO 11898 · AUTOSAR SecOC · ISO 21434",
                 Inches(1), Inches(5.5), Inches(11.33), Inches(1.0),
                 font_size=12, colour=C_LIGHT_GREY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Main assembler
# ─────────────────────────────────────────────────────────────────────────────

def build_presentation(output_path: str = "Sentinel_T_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides …")
    slide_title(prs);              print("  ✔  Slide 1: Title")
    slide_agenda(prs);             print("  ✔  Slide 2: Agenda")
    slide_intro(prs);              print("  ✔  Slide 3: Introduction")
    slide_existing_systems(prs);   print("  ✔  Slide 4: Existing Systems")
    slide_proposed_solution(prs);  print("  ✔  Slide 5: Proposed Solution")
    slide_architecture(prs);       print("  ✔  Slide 6: Architecture")
    slide_results_live(prs);       print("  ✔  Slide 7: Live Results")
    slide_results_datasets(prs);   print("  ✔  Slide 8: Dataset Results")
    slide_results_stress(prs);     print("  ✔  Slide 9: Stress Test")
    slide_conclusion(prs);         print("  ✔  Slide 10: Conclusion")
    slide_thankyou(prs);           print("  ✔  Slide 11: Thank You")

    prs.save(output_path)
    print(f"\n✅  Presentation saved → {output_path}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build_presentation()
